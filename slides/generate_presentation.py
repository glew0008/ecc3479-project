from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Education and Income in Australia'
slide.placeholders[1].text = 'ABS Census 2016 and 2021 | Australian males'

content = [
    ('Research question and why it matters',
     'How does educational attainment affect personal income for Australian males, and how did that relationship change from 2016 to 2021?\n\nWhy it matters:\n- Education-income gradients are central to inequality and labour-market policy.\n- Governments, students, and employers all rely on returns-to-education evidence.\n- The census data let us compare population-level outcomes across two Census years.'),
    ('Data',
     'Source: ABS Census TableBuilder 2016 and 2021.\nSample: Australian males only; 8 substantive education groups.\nExclusions: Total, Not stated, Not applicable, Supplementary Codes.\nKey variables: education group, income bracket, cell counts, year.\nData treatment: harmonised top income brackets and converted bracket values to weekly midpoints.'),
    ('Empirical strategy',
     'Model: weighted mean income = α + β * education rank + γ * year indicator.\nKey coefficient: β is the average weekly income gap per education rank step.\nWeights: use cell counts because rows represent different population sizes.\nAssumption: this is descriptive, not causal. Education rank is not exogenous to income because age, occupation, hours, and location may differ by education group.'),
    ('Main results',
     'Headline finding: each one-step increase in education rank is associated with about $196 more per week in weighted mean income.\nControlling for year, 2021 incomes are also about $121/week higher than 2016.\nThe low-to-high education gap is roughly $1,372/week, or about $71,300/year.\nThe pattern is visible in every year and is consistent across main groups.'),
    ('Robustness and limitations',
     'Robustness: the association survives unweighted models, no-year models, alternative income transformations, and year-specific samples.\nLimitation: grouped census data lack individual-level controls for age, occupation, hours worked, region, and firm characteristics.\nLargest threat: omitted-variable bias from composition differences across education groups and measurement error from income bracket midpoints.'),
    ('Conclusion and next steps',
     'Conclusion: there is a strong descriptive education-income gradient for Australian males in 2016 and 2021 census aggregates.\nWhat we cannot conclude: a causal return to education from this grouped data.\nNext step: use individual-level microdata with controls, or seek quasi-experimental variation or longitudinal data to address the main threat.'),
]

for title, text in content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.text = text
    for paragraph in body.paragraphs:
        paragraph.font.size = Pt(18)

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Main results'
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.8))
text_frame = textbox.text_frame
text_frame.text = 'Headline finding: each one-step increase in education rank is associated with about $196 more per week in weighted mean income.'
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(18)

img_path = Path('outputs/eda/figures/weighted_mean_income_by_education.png')
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(3.2), width=Inches(9))
else:
    print(f'Warning: image {img_path} not found')

slides_dir = Path('slides')
slides_dir.mkdir(parents=True, exist_ok=True)
prs.save(slides_dir / 'Education_and_Income_Presentation.pptx')
print('Created slides/Education_and_Income_Presentation.pptx')
